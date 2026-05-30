from __future__ import annotations

import asyncio
import csv
import json
import re
import shlex
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import List

from luoying_bot.domain.context import ChatContext
from luoying_bot.ports.transport import ChatTransport

MAX_READ_FILE_CHARS = 120_000

@dataclass(slots=True)
class ScriptOpResult:
    ok: bool
    text: str
    data: dict

class ScriptWorkspaceService:
    def __init__(
        self,
        root_dir:Path,
        python_timeout_sec:int=15,
    ):
        self.root_dir=Path(root_dir)
        self.root_dir.mkdir(parents=True,exist_ok=True)
        self.python_timeout_sec=python_timeout_sec
    
    def _user_dir(self,user_id:str)->Path:
        path=self.root_dir/str(user_id)
        path.mkdir(parents=True,exist_ok=True)
        return path
    
    def _sanitive_relative_path(self, file_path:str)->str:
        raw=(file_path or "").strip().replace("\\",'/')
        if not raw:
            raise ValueError("file_path 不能为空")
        if raw.startswith("/"):
            raise ValueError("不允许绝对路径")
        parts = [p for p in raw.split("/") if p not in ("", ".")]
        if not parts:
            raise ValueError("非法文件路径")
        if any(part == ".." for part in parts):
            raise ValueError("不允许使用 ..")
        return "/".join(parts)
    
    def _resolve_user_file(self,user_id:str,file_path:str)->Path:
        rel=self._sanitive_relative_path(file_path)
        base=self._user_dir(user_id=user_id).resolve()
        target=(base/rel).resolve()
        if base != target and base not in target.parents:
            raise ValueError("文件路径越界")
        return target

    def _read_text_smart(self, target: Path) -> tuple[str, str]:
        raw = target.read_bytes()
        for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk", "big5"):
            try:
                return raw.decode(encoding), encoding
            except UnicodeDecodeError:
                continue
        return raw.decode("latin1"), "latin1"

    def _trim_read_content(self, content: str) -> str:
        content = content.replace("\x00", "").strip()
        if len(content) <= MAX_READ_FILE_CHARS:
            return content
        return (
            content[:MAX_READ_FILE_CHARS]
            + f"\n\n[内容过长，已截断到前 {MAX_READ_FILE_CHARS} 个字符]"
        )

    def _is_probably_binary(self, raw: bytes) -> bool:
        if b"\x00" in raw[:4096]:
            return True
        if not raw:
            return False
        sample = raw[:4096]
        control = sum(1 for byte in sample if byte < 32 and byte not in (9, 10, 13))
        return control / len(sample) > 0.08

    def _extract_binary_text_runs(self, raw: bytes) -> str:
        ascii_text = re.findall(rb"[ -~]{4,}", raw)
        utf16_text = re.findall(rb"(?:[ -~]\x00){4,}", raw)
        parts: list[str] = []
        seen: set[str] = set()

        for item in ascii_text:
            text = item.decode("latin1", errors="ignore").strip()
            if text and text not in seen:
                seen.add(text)
                parts.append(text)

        for item in utf16_text:
            text = item.decode("utf-16le", errors="ignore").strip()
            if text and text not in seen:
                seen.add(text)
                parts.append(text)

        return self._trim_read_content("\n".join(parts))

    def _read_pdf_text(self, target: Path) -> str:
        from pypdf import PdfReader

        reader = PdfReader(str(target))
        pages: list[str] = []
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            pages.append(f"[第 {index} 页]\n{text}".strip())
        return "\n\n".join(page for page in pages if page)

    def _read_xlsx_text(self, target: Path) -> str:
        from openpyxl import load_workbook

        workbook = load_workbook(target, read_only=True, data_only=True)
        parts: list[str] = []
        try:
            for sheet in workbook.worksheets:
                parts.append(f"[工作表：{sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    values = ["" if value is None else str(value) for value in row]
                    if any(value.strip() for value in values):
                        parts.append("\t".join(values).rstrip())
        finally:
            workbook.close()
        return "\n".join(parts)

    def _read_xls_text(self, target: Path) -> str:
        import xlrd

        workbook = xlrd.open_workbook(str(target))
        parts: list[str] = []
        for sheet in workbook.sheets():
            parts.append(f"[工作表：{sheet.name}]")
            for row_index in range(sheet.nrows):
                values = [str(sheet.cell_value(row_index, col_index)) for col_index in range(sheet.ncols)]
                if any(value.strip() for value in values):
                    parts.append("\t".join(values).rstrip())
        return "\n".join(parts)

    def _read_docx_text(self, target: Path) -> str:
        from docx import Document

        document = Document(str(target))
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells]
                if any(values):
                    parts.append("\t".join(values))
        return "\n".join(parts)

    def _read_pptx_text(self, target: Path) -> str:
        from pptx import Presentation

        presentation = Presentation(str(target))
        parts: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                parts.append(f"[幻灯片 {index}]\n" + "\n".join(slide_parts))
        return "\n\n".join(parts)

    def _read_csv_text(self, target: Path, delimiter: str | None = None) -> str:
        content, encoding = self._read_text_smart(target)
        dialect = csv.excel_tab if delimiter == "\t" else csv.excel
        rows = csv.reader(content.splitlines(), dialect=dialect)
        return "\n".join("\t".join(cell for cell in row) for row in rows)

    def _read_html_text(self, target: Path) -> str:
        from bs4 import BeautifulSoup

        content, _ = self._read_text_smart(target)
        soup = BeautifulSoup(content, "html.parser")
        return soup.get_text("\n")

    def _read_json_text(self, target: Path) -> str:
        content, _ = self._read_text_smart(target)
        return json.dumps(json.loads(content), ensure_ascii=False, indent=2)

    def _read_rtf_text(self, target: Path) -> str:
        content, _ = self._read_text_smart(target)
        content = re.sub(r"\\'[0-9a-fA-F]{2}", " ", content)
        content = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", content)
        content = content.replace("{", " ").replace("}", " ")
        content = content.replace("\\", " ")
        return re.sub(r"\s+", " ", content).strip()

    def _read_odf_text(self, target: Path) -> str:
        from odf.opendocument import load
        from odf import table, text
        from odf.teletype import extractText

        document = load(str(target))
        parts: list[str] = []
        for paragraph in document.getElementsByType(text.P):
            value = extractText(paragraph).strip()
            if value:
                parts.append(value)
        for cell in document.getElementsByType(table.TableCell):
            value = extractText(cell).strip()
            if value:
                parts.append(value)
        return "\n".join(parts)

    def _read_with_libreoffice_text(self, target: Path) -> str | None:
        exe = shutil.which("soffice") or shutil.which("libreoffice")
        if not exe:
            return None
        with tempfile.TemporaryDirectory() as tmp:
            proc = subprocess.run(
                [
                    exe,
                    "--headless",
                    "--convert-to",
                    "txt:Text",
                    "--outdir",
                    tmp,
                    str(target),
                ],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=45,
                check=False,
            )
            if proc.returncode != 0:
                return None
            outputs = list(Path(tmp).glob("*.txt"))
            if not outputs:
                return None
            content, _ = self._read_text_smart(outputs[0])
            return content

    def _read_structured_text(self, target: Path) -> tuple[str, str] | None:
        suffix = target.suffix.lower()
        if suffix == ".pdf":
            return self._read_pdf_text(target), "PDF 文本提取"
        if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
            return self._read_xlsx_text(target), "Excel 文本提取"
        if suffix == ".xls":
            return self._read_xls_text(target), "Excel 文本提取"
        if suffix == ".docx":
            return self._read_docx_text(target), "Word 文本提取"
        if suffix == ".pptx":
            return self._read_pptx_text(target), "PowerPoint 文本提取"
        if suffix == ".csv":
            return self._read_csv_text(target), "CSV 文本提取"
        if suffix == ".tsv":
            return self._read_csv_text(target, delimiter="\t"), "TSV 文本提取"
        if suffix in {".html", ".htm", ".xml"}:
            return self._read_html_text(target), "HTML/XML 文本提取"
        if suffix == ".json":
            return self._read_json_text(target), "JSON 格式化读取"
        if suffix == ".rtf":
            return self._read_rtf_text(target), "RTF 文本提取"
        if suffix in {".odt", ".ods", ".odp"}:
            return self._read_odf_text(target), "OpenDocument 文本提取"
        if suffix in {".doc", ".ppt"}:
            libreoffice_text = self._read_with_libreoffice_text(target)
            if libreoffice_text is not None:
                return libreoffice_text, "LibreOffice 文本提取"
        return None

    def list_scripts(self,user_id:str)->ScriptOpResult:
        base = self._user_dir(user_id=user_id)
        files = [p for p in base.rglob("*") if p.is_file()]
        rels = [
            str(p.relative_to(base)).replace("\\", "/")
            for p in files
        ]
        rels.sort()

        if not rels:
            return ScriptOpResult(True,"当前没有脚本文件",{"files":[]})
        lines = ["当前脚本文件："]
        lines.extend(f"{i}. {name}" for i, name in enumerate(rels, start=1))
        return ScriptOpResult(True,"\n".join(lines), {"files": rels})

    def tree(self, user_id: str) -> ScriptOpResult:
        base = self._user_dir(user_id=user_id)

        def visible_children(path: Path) -> list[Path]:
            return sorted(
                (child for child in path.iterdir() if not child.name.startswith(".")),
                key=lambda child: (not child.is_dir(), child.name.lower()),
            )

        def render(path: Path, prefix: str = "") -> list[str]:
            lines: list[str] = []
            children = visible_children(path)
            for index, child in enumerate(children):
                is_last = index == len(children) - 1
                connector = "└── " if is_last else "├── "
                name = f"{child.name}/" if child.is_dir() else child.name
                lines.append(f"{prefix}{connector}{name}")
                if child.is_dir():
                    extension = "    " if is_last else "│   "
                    lines.extend(render(child, prefix + extension))
            return lines

        lines = [f"{base.name}/"]
        lines.extend(render(base))
        if len(lines) == 1:
            lines.append("(空)")
        return ScriptOpResult(True, "\n".join(lines), {"root": str(base)})

    def read_script(self,user_id:str,file_path:str)->ScriptOpResult:
        target = self._resolve_user_file(user_id=user_id,file_path=file_path)
        if not target.exists() or not target.is_file():
            return ScriptOpResult(False, f"文件不存在：{file_path}", {"file_path": file_path})

        structured_error: str | None = None
        try:
            structured = self._read_structured_text(target)
        except Exception as exc:
            structured = None
            structured_error = f"{type(exc).__name__}: {exc}"
        if structured is not None:
            content, reader = structured
            content = self._trim_read_content(content or "(未提取到文本)")
            return ScriptOpResult(
                True,
                f"文件内容如下（{reader}）：\n{content}",
                {
                    "file_path": str(target),
                    "content": content,
                    "reader": reader,
                    "structured": True,
                },
            )

        raw = target.read_bytes()
        if self._is_probably_binary(raw):
            content = self._extract_binary_text_runs(raw)
            if content:
                note = "结构化解析失败，已提取二进制文件中的可见文本片段"
                if structured_error:
                    note += f"；结构化解析错误：{structured_error}"
                return ScriptOpResult(
                    True,
                    f"文件内容如下（{note}）：\n{content}",
                    {
                        "file_path": str(target),
                        "content": content,
                        "reader": "binary_strings",
                        "structured": False,
                        "structured_error": structured_error,
                    },
                )
            return ScriptOpResult(
                False,
                f"无法读取该二进制文件中的文本内容：{file_path}"
                + (f"\n结构化解析错误：{structured_error}" if structured_error else ""),
                {"file_path": str(target), "structured_error": structured_error},
            )

        content, encoding = self._read_text_smart(target)
        content = self._trim_read_content(content)
        return ScriptOpResult(
            True,
            f"文件内容如下（自动识别编码：{encoding}）：\n{content}",
            {
                "file_path": str(target),
                "content": content,
                "encoding": encoding,
                "structured": False,
                "structured_error": structured_error,
            },
        )

    def write_script(self, user_id: str, file_path: str, content: str, overwrite: bool = False) -> ScriptOpResult:
        target = self._resolve_user_file(user_id,file_path)
        target.parent.mkdir(parents=True, exist_ok=True)

        if target.exists() and not overwrite:
            return ScriptOpResult(
                False,
                f"文件已存在：{file_path}，如需覆盖请调用 overwrite",
                {"file_path": file_path},
            )
        
        target.write_text(content or "",encoding="utf-8")
        action = "已覆盖" if overwrite else "已创建"
        return ScriptOpResult(
            True,
            f"{action}脚本：{file_path}",
            {
                "file_path": str(target),
                "size": len(content or ""),
                "overwrite": overwrite,
            },
        )
    
    def delete_script(self, user_id: str, file_path: str) -> ScriptOpResult:
        target = self._resolve_user_file(user_id, file_path)
        if not target.exists() or not target.is_file():
            return ScriptOpResult(False, f"文件不存在：{file_path}", {"file_path": file_path})
        target.unlink()
        return ScriptOpResult(True, f"已删除脚本：{file_path}", {"file_path": file_path})

    async def run_python_script(self, user_id: str, file_path: str, args: str = "") -> ScriptOpResult:
        target = self._resolve_user_file(user_id, file_path)
        if not target.exists() or not target.is_file():
            return ScriptOpResult(False, f"文件不存在：{file_path}", {"file_path": file_path})

        if target.suffix.lower() != ".py":
            return ScriptOpResult(False, "只能运行 .py 文件", {"file_path": file_path})

        argv: List[str] = shlex.split(args or "")
        proc = await asyncio.create_subprocess_exec(
            sys.executable,
            str(target),
            *argv,
            cwd=str(target.parent),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )

        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=self.python_timeout_sec)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.communicate()
            out = ""
            err = f"脚本运行超时：{self.python_timeout_sec}s"
            combined = (
                f"script: {file_path}\n"
                f"args: {args or '(none)'}\n"
                f"exit_code: timeout\n\n"
                f"[stdout]\n(timeout before capture completed)\n\n"
                f"[stderr]\n{err}"
            )

            return ScriptOpResult(
                False,
                combined,
                {
                    "type": "script_result",
                    "file_path": file_path,
                    "args": args or "",
                    "returncode": None,
                    "timeout": self.python_timeout_sec,
                    "stdout": out,
                    "stderr": err,
                },
            )

        out = (stdout or b"").decode("utf-8", errors="replace")
        err = (stderr or b"").decode("utf-8", errors="replace")
        full_combined = (
            f"script: {file_path}\n"
            f"args: {args or '(none)'}\n"
            f"exit_code: {proc.returncode}\n\n"
            f"[stdout]\n{out or '(empty)'}\n\n"
            f"[stderr]\n{err or '(empty)'}"
        )
        return ScriptOpResult(
            proc.returncode == 0,
            full_combined,
            {
                "type": "script_result",
                "file_path": file_path,
                "args": args or "",
                "returncode": proc.returncode,
                "stdout": out,
                "stderr": err,
                "timeout": False,
            },
        )

    async def send_script_to_transport(
        self,
        user_id: str,
        file_path: str,
        context: ChatContext,
        transport: ChatTransport,
    ) -> ScriptOpResult:
        target = self._resolve_user_file(user_id, file_path)
        if not target.exists() or not target.is_file():
            return ScriptOpResult(False, f"文件不存在：{file_path}", {"file_path": file_path})

        await transport.upload_file(context=context, file=str(target))

        return ScriptOpResult(
            True,
            f"已发送脚本到当前会话：{file_path}",
            {"file_path": file_path},
        )
